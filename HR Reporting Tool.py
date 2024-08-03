# Report Automation Project

# Author: Oscar García
# Purpose: automating the periodic creation of a PowerPoint report based on sql relational database Employees, available
# under https://downloads.mysql.com/docs/employee-en.pdf and https://dev.mysql.com/doc/index-other.html

# 0. Preparatory tasks
# Press Shift+F10 to execute it or replace it with your code.
# Press Double Shift to search everywhere for classes, files, tool windows, actions, and settings.

# 0.a. Load packages
import pandas as pd
import os
import matplotlib.pyplot as plt
import mysql.connector
import seaborn as sns
import numpy as np
import statsmodels.api as sm
from sklearn.linear_model import LinearRegression
from statsmodels.formula.api import mixedlm
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor  # Import RGBColor for setting colors
from pptx.enum.text import PP_ALIGN
from pandas.tseries.offsets import DateOffset

# from pptx.chart.data import CategoryChartData
# from pptx.enum.chart import XL_CHART_TYPE

# 0.b. Prepare sql connection
# Database connection
db_config = {
    'user': 'root',
    'password': '250598',
    'host': 'localhost',
    'database': 'employees',
    'charset': 'utf8mb4',
    'collation': 'utf8mb4_general_ci'
}

# Connect to the database
cnx = mysql.connector.connect(**db_config)
cursor = cnx.cursor()

# 0.c. Set up output directory for plots and slides
# Directory to save the plots
output_dir = r'C:\Users\oscar\Documents\Projects\Report Automation Project\PythonOutput'

# Ensure the directory exists
if not os.path.exists(output_dir):
    os.makedirs(output_dir)

# Debug print to ensure the directory is correct
print(f"Output directory: {output_dir}")

# 0.d. Set up slides generation
prs = Presentation()


# Function to set the footer for each slide
def set_footer(slide, slide_number):
    left_text = "CONFIDENTIAL - Do not distribute"
    right_text = datetime.now().strftime("%B %d, 2024")

    # Add left text box
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(5), Inches(0.5))
    left_frame = left_box.text_frame
    left_frame.text = left_text
    for paragraph in left_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.name = 'Segoe UI Light'
        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black color

    # Add right text box
    right_box = slide.shapes.add_textbox(Inches(8), Inches(7), Inches(2), Inches(0.5))
    right_frame = right_box.text_frame
    right_frame.text = f"{right_text} - Pg. {slide_number}"
    for paragraph in right_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.name = 'Segoe UI Light'
        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black color


# Function to set background color for each slide
def set_background(slide, color=RGBColor(255, 255, 255), img_path=None):  # Light grey by default
    # Light grey by default background
    background = slide.background
    fill = background.fill

    # Set background to solid color if no image is provided
    if not img_path:
        fill.solid()
        fill.fore_color.rgb = color
    else:
        left = top = Inches(0)
        img_path = img_path
        pic = slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)

        # This moves it to the background
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)


# Create slides with desired content
def create_slide(prs, title_text, content=None, img_path=None):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text

    # Set title font to Georgia, size to 24, and align left
    title_frame = title.text_frame
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.name = 'Segoe UI Semibold'
    title_paragraph.font.size = Pt(30)
    title_paragraph.font.color.rgb = RGBColor(23, 55, 94)
    title_paragraph.alignment = PP_ALIGN.LEFT

    # Define text box position and size
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(5.0)

    if content:
        # Add a textbox for the content
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(5.5))
        text_frame = textbox.text_frame
        for item in content:
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(24)
            p.font.name = 'Segoe UI Light'

    # Add picture if img_path is provided
    if img_path:
        slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8.5))


    slide_number = len(prs.slides)
    set_footer(slide, slide_number)
    set_background(slide, img_path='C:/Users/oscar/Documents/Projects/Report Automation Project/Supporting materials/slide_background2.jpg')


# 0.e. Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Employees Report"
subtitle.text = "Automated Report Generated with Python"

# Set title and subtitle font to Georgia
title.text_frame.paragraphs[0].font.name = 'Segoe UI Semibold'
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(23, 55, 94)
subtitle.text_frame.paragraphs[0].font.name = 'Segoe UI Light'

set_footer(slide, 1)
set_background(slide,
               img_path='C:/Users/oscar/Documents/Projects/Report Automation Project/Supporting materials/slide_background2.jpg')  # Light grey background

# 0.f. Slide 2: Index
# Create slides with desired content
def create_slide_index(prs, title_text, content=None):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text

    # Set title font to Georgia, size to 24, and align left
    title_frame = title.text_frame
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.name = 'Segoe UI Semibold'
    title_paragraph.font.size = Pt(30)
    title_paragraph.font.color.rgb = RGBColor(23, 55, 94)
    title_paragraph.alignment = PP_ALIGN.LEFT

    # Define text box position and size
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(5.0)

    if content:
        # Add a textbox for the content
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(5.5))
        text_frame = textbox.text_frame
        for item in content:
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.font.name = 'Segoe UI Light'

    slide_number = len(prs.slides)
    set_footer(slide, slide_number)
    set_background(slide,
                   img_path='C:/Users/oscar/Documents/Projects/Report Automation Project/Supporting materials/slide_background2.jpg')


content = [
    f"Current Employees Overview..............................................................................3",
    f"Number of Employees and Average Salary (USD, K) by Department......4",
    f"Number of Employees and Average Salary (USD, K) by Title.....................5",
    f"ESRS Disclosure requirement S1-7: Gender Ratio across Workforce........6",
    f"ESRS Disclosure requirement S1-7: Gender Ratio by Title...........................7",
    f"ESRS Disclosure Requirement S1-7: Gender Ratio by Department...........8",
    f"ESRS Disclosure requirement S1-16: Gender Pay Gap..................................9",
    f"ESRS Disclosure requirement S1-16: Gender Pay Gap by Title..................10",
    f"ESRS Disclosure requirement S1-16: Gender Pay Gap by Department...11",
    f"ESRS Disclosure requirement S1-16: Gender Pay Gap over Time............12",
    f"Number of Employees Over Time...................................................................13",
    f"ESRS S1-7: Gender Workforce Composition Over Time............................14",
    f"Workforce Composition by Department Over Time..................................15",
    f"Salary Progression by Gender..........................................................................16",
    f"Average Starting Salary by Gender Over Time...........................................17"
]

create_slide_index(prs, "Index of Content", content=content)

# 1. Overview of HR-relevant data
# 1.a. Fetching & processing data
# Define the SQL query with multiple JOINs and condition
query = """
SELECT e.emp_no, e.gender, e.birth_date, e.hire_date, n.dept_name, t.title, s.salary
FROM employees e
JOIN dept_emp d ON e.emp_no = d.emp_no AND d.to_date = '9999-01-01'
JOIN departments n ON d.dept_no = n.dept_no
JOIN titles t ON e.emp_no = t.emp_no AND t.to_date = '9999-01-01'
JOIN salaries s ON e.emp_no = s.emp_no AND s.to_date = '9999-01-01'
WHERE d.to_date = '9999-01-01';
"""
# Execute the query
cursor.execute(query)
# Fetch all the results
data = cursor.fetchall()
curr_df = pd.DataFrame(data, columns=['emp_no', 'gender', 'birth_date', 'hire_date', 'dept_name', 'title', 'salary'])
# Print the results
print(curr_df.head(15))

# 1.b. Analysis
# Calculate required statistics
num_employees = curr_df.shape[0]
average_salary = curr_df['salary'].mean()
std_salary = curr_df['salary'].std()
total_annual_salary_cost = curr_df['salary'].sum()

# 1.c. Slide generation
# Slide 2: Current Employees Analysis
content = [
    f"Number of Employees: {num_employees}",
    f"",
    f"Average Salary: ${average_salary:,.2f}",
    f"",
    f"Standard Deviation of Salary: ${std_salary:,.2f}",
    f"",
    f"Total Annual Salary Cost: ${total_annual_salary_cost:,.2f}"
]

create_slide(prs, "Current Employees Overview", content=content)

# 2. Distributions by department and title
# 2.a. Department - Calculate the necessary statistics
dept_stats = curr_df.groupby('dept_name').agg(
    num_employees=('emp_no', 'nunique'),
    avg_salary=('salary', 'mean')
).reset_index()

# Divide the salary by 1000
dept_stats['avg_salary'] = dept_stats['avg_salary'] / 1000

# 2.b. Department - Plot the data
fig, ax1 = plt.subplots(figsize=(12, 6))

# Bar plot for frst y-axis number of employees
color = '#3179B5'
ax1.set_xlabel('Department')
ax1.set_ylabel('Number of Employees', color=color)
bars = ax1.bar(dept_stats['dept_name'], dept_stats['num_employees'], color=color, alpha=0.7,
               label='Number of Employees')
ax1.tick_params(axis='y', labelcolor=color)

# Add text labels for number of employees
for bar in bars:
    height = bar.get_height()
    ax1.text(bar.get_x() + bar.get_width() / 2, height - (height * 0.05), f'{height}', ha='center', va='top',
             color='white', fontsize=12, fontweight='bold')

# Create another y-axis for the average salary
ax2 = ax1.twinx()
color = '#AEDCF8'
ax2.set_ylabel('Average Salary (in thousands)', color=color)
line, = ax2.plot(dept_stats['dept_name'], dept_stats['avg_salary'], color=color, marker='o', linestyle='-',
                 linewidth=2, label='Average Salary')
ax2.tick_params(axis='y', labelcolor=color)

# Add text labels for average salary
for i, avg_salary in enumerate(dept_stats['avg_salary']):
    ax2.text(i, avg_salary, f'{avg_salary:.1f}', ha='center', va='bottom', color='#104B66', fontsize=12,
             fontweight='bold')

# Add titles and legend
fig.suptitle('Number of Employees and Average Salary (USD, K) by Department')
fig.tight_layout()  # Adjust layout to make room for the labels
# Rotate x-axis labels for better readability
plt.xticks(rotation=45, ha='right')
plt.savefig(f'{output_dir}/num_emp_sal_per_dept.png')
plt.close()

# 2.c. Department - Generate slides
# Slide 3: Line Plot
create_slide(prs, "Number of Employees and Average Salary (USD, K) by Department",
             img_path=(f'{output_dir}\\num_emp_sal_per_dept.png'))


# 2.d. Title - Calculate the necessary statistics
dept_stats = curr_df.groupby('title').agg(
    num_employees=('emp_no', 'nunique'),
    avg_salary=('salary', 'mean')
).reset_index()

# Divide the salary by 1000
dept_stats['avg_salary'] = dept_stats['avg_salary'] / 1000

# 2.e. Title - Plot the data
fig, ax1 = plt.subplots(figsize=(12, 6))

# Bar plot for frst y-axis number of employees
color = '#3179B5'
ax1.set_xlabel('Title')
ax1.set_ylabel('Number of Employees', color=color)
bars = ax1.bar(dept_stats['title'], dept_stats['num_employees'], color=color, alpha=0.7, label='Number of Employees')
ax1.tick_params(axis='y', labelcolor=color)

# Add text labels for number of employees
for bar in bars:
    height = bar.get_height()
    ax1.text(bar.get_x() + bar.get_width() / 2, height - (height * 0.05), f'{height}', ha='center', va='top',
             color='white', fontsize=12, fontweight='bold')

# Create another y-axis for the average salary
ax2 = ax1.twinx()
color = '#AEDCF8'
ax2.set_ylabel('Average Salary (in thousands)', color=color)
line, = ax2.plot(dept_stats['title'], dept_stats['avg_salary'], color=color, marker='o', linestyle='-', linewidth=2,
                 label='Average Salary')
ax2.tick_params(axis='y', labelcolor=color)

# Add text labels for average salary
for i, avg_salary in enumerate(dept_stats['avg_salary']):
    ax2.text(i,
             avg_salary, f'{avg_salary:.1f}', ha='center', va='bottom', color='#104B66', fontsize=12, fontweight='bold')

# Add titles and legend
fig.suptitle('Number of Employees and Average Salary by Title')
fig.tight_layout()  # Adjust layout to make room for the labels
# Rotate x-axis labels for better readability
plt.xticks(rotation=45, ha='right')

# Save the plot
plt.savefig(f'{output_dir}/num_emp_sal_per_titl.png')
plt.close()

# 2.f. Title - Generate slides
# Slide 4: Line Plot
create_slide(prs, "Number of Employees and Average Salary (USD, K) by Title",
             img_path=f'{output_dir}\\num_emp_sal_per_titl.png')

# 3. Gender gap measures.
# ESRS S1. Own workforce, Equal treatment and opportunities for all.
# Gender equality and equal pay for work of equal value
# Disclosure requirement S1-7  – Characteristics of the Undertaking’s Employees
# AG 90. Providing a breakdown of employees by gender and type of employment relationship gives
# insight into gender representation across the undertaking.
# 3.a. Fetching & processing data
# Define the SQL query with multiple JOINs and condition
query = """
SELECT e.emp_no, e.gender, e.birth_date, n.dept_name, t.title, s.salary
FROM employees e
JOIN dept_manager d ON e.emp_no = d.emp_no AND d.to_date = '9999-01-01'
JOIN departments n ON d.dept_no = n.dept_no
JOIN titles t ON e.emp_no = t.emp_no AND t.to_date = '9999-01-01'
JOIN salaries s ON e.emp_no = s.emp_no AND s.to_date = '9999-01-01'
WHERE d.to_date = '9999-01-01';
"""
# Execute the query
cursor.execute(query)
# Fetch all the results
data = cursor.fetchall()
manager_df = pd.DataFrame(data, columns=['emp_no', 'gender', 'birth_date', 'dept_name', 'title', 'salary'])
# Print the results
print(manager_df.head(20))

# Calculate the number and percentage of managers by gender
manager_gender_counts = manager_df['gender'].value_counts()

# Prepare data for the manager pie chart
manager_labels = [f'{gender} ({count})' for gender, count in manager_gender_counts.items()]
manager_sizes = manager_gender_counts.values

# Count the number of male and female employees in the current DataFrame
gender_counts = curr_df['gender'].value_counts()

# Prepare data for the current employee pie chart
labels = [f'{gender} ({count})' for gender, count in gender_counts.items()]
sizes = gender_counts.values

# Create a figure with two subplots side by side
fig, axes = plt.subplots(1, 2, figsize=(18, 11))

# Plot the pie chart for managers
axes[0].pie(manager_sizes, labels=manager_labels, autopct='%1.1f%%', startangle=140, colors=['#66b3ff','#ff9999'], textprops={'fontsize': 24})
axes[0].set_title('Top Management’s Gender Composition', fontsize=24)
axes[0].axis('equal')  # Equal aspect ratio ensures that pie is drawn as a circle

# Plot the pie chart for current employees
axes[1].pie(sizes, labels=labels, autopct='%1.1f%%', startangle=140, colors=['#66b3ff','#ff9999'], textprops={'fontsize': 24})
axes[1].set_title('Total Employees’ Gender Composition', fontsize=24)
axes[1].axis('equal')  # Equal aspect ratio ensures that pie is drawn as a circle

# Adjust layout to ensure everything fits without overlapping
plt.tight_layout()

# Save the combined plot
plt.savefig(f'{output_dir}/combined_gender_distribution.png')
plt.close()

create_slide(prs, "ESRS Disclosure requirement S1-7: Gender Ratio across Workforce", img_path=f'{output_dir}\\combined_gender_distribution.png')

# Group by title and gender to get the counts
title_gender_counts = curr_df.groupby(['title', 'gender']).size().unstack(fill_value=0)

# Plot the column chart
title_gender_counts.plot(kind='bar', figsize=(14, 7), color=['#ff9999', '#66b3ff'])

# Adding labels and title
plt.xlabel('Job Title')
plt.ylabel('Number of Employees')
plt.title('Distribution of Employees by Title and Gender')

# Adding legend
plt.legend(title='Gender')

# Display the plot
plt.xticks(rotation=45, ha='right')
plt.tight_layout()
plt.savefig(f'{output_dir}/gender_ratio_by_title.png')
plt.close()

create_slide(prs, "ESRS Disclosure requirement S1-7: Gender Ratio by Title", img_path=f'{output_dir}\\gender_ratio_by_title.png')

# Group by department name and gender to get the counts
dept_gender_counts = curr_df.groupby(['dept_name', 'gender']).size().unstack(fill_value=0)

# Plot the column chart
dept_gender_counts.plot(kind='bar', figsize=(14, 7), color=['#ff9999', '#66b3ff'])

# Adding labels and title
plt.xlabel('Department Name')
plt.ylabel('Number of Employees')
plt.title('Distribution of Employees by Department and Gender')

# Adding legend
plt.legend(title='Gender')

# Display the plot
plt.xticks(rotation=45, ha='right')
plt.tight_layout()
plt.savefig(f'{output_dir}/gender_ratio_by_dept.png')
plt.close()

create_slide(prs, "ESRS Disclosure Requirement S1-7: Gender Ratio by Department", img_path=f'{output_dir}\\gender_ratio_by_dept.png')

# Disclosure requirement S1-16 – Pay gap between women and men
# The undertaking shall disclose the following information:
# (a) the male-female pay gap, defined as the difference between average gross hourly
# earnings of male paid employees and of female paid employees expressed as a
# percentage of average gross hourly earnings of male paid employees;


# 3.a. Analysing current data
def calculate_pay_gap_for_date(df):
    avg_salary_male = df[df['gender'] == 'M']['salary'].mean()
    avg_salary_female = df[df['gender'] == 'F']['salary'].mean()

    if avg_salary_male > 0:
        pay_gap = ((avg_salary_male - avg_salary_female) / avg_salary_male) * 100
    else:
        pay_gap = 0  # Avoid division by zero

    return pay_gap


# Calculate the pay gap for a specific day in the current year
current_pay_gap = calculate_pay_gap_for_date(curr_df)
print(current_pay_gap)

content = [
    f"The gender pay gap, calculated as the difference between",
    f"average gross hourly earnings of male paid employees and",
    f"of female paid employees expressed as a percentage of",
    f"average gross hourly earnings of male paid employees, is:",
    f"",
    f" {current_pay_gap:.2f}%"
]

# Calculate the pay gap by title
pay_gap_by_title = curr_df.groupby('title').apply(calculate_pay_gap_for_date).reset_index(name='pay_gap')

# Calculate the pay gap by department
pay_gap_by_dept = curr_df.groupby('dept_name').apply(calculate_pay_gap_for_date).reset_index(name='pay_gap')

# Plot the pay gap by title
plt.figure(figsize=(12, 6))
bars = plt.bar(pay_gap_by_title['title'], pay_gap_by_title['pay_gap'], color='skyblue')
plt.xlabel('Title')
plt.ylabel('Pay Gap (%)')
plt.title('Gender Pay Gap by Title')
plt.xticks(rotation=90)

# Add text labels
for bar in bars:
    height = bar.get_height()
    plt.text(bar.get_x() + bar.get_width() / 2, height, f"{height:.2f}%", ha='center', va='bottom', fontsize=9)

plt.tight_layout()
plt.savefig(f'{output_dir}/gender_paygap_by_title.png')
plt.close()

# Plot the pay gap by department
plt.figure(figsize=(12, 6))
bars = plt.bar(pay_gap_by_dept['dept_name'], pay_gap_by_dept['pay_gap'], color='lightcoral')
plt.xlabel('Department')
plt.ylabel('Pay Gap (%)')
plt.title('Gender Pay Gap by Department')
plt.xticks(rotation=90)

# Add text labels
for bar in bars:
    height = bar.get_height()
    plt.text(bar.get_x() + bar.get_width() / 2, height, f"{height:.2f}%", ha='center', va='bottom', fontsize=9)

plt.tight_layout()
plt.savefig(f'{output_dir}/gender_paygap_by_dept.png')
plt.close()

# 3.b. Fetching past data
years_query = """
SELECT DISTINCT YEAR(from_date) AS year
FROM dept_emp
UNION
SELECT DISTINCT YEAR(from_date) AS year
FROM titles
UNION
SELECT DISTINCT YEAR(from_date) AS year
FROM salaries;
"""
cursor.execute(years_query)
years_data = cursor.fetchall()
years = [row[0] for row in years_data]


# Function to calculate the number of days in a given year
def days_in_year(year):
    return (datetime(year + 1, 1, 1) - datetime(year, 1, 1)).days


# DataFrame to store annual salaries
df_annsal = pd.DataFrame()

# Iterate over each year and execute the modified query
for year in years:
    start_date = datetime(year, 1, 1)
    end_date = datetime(year, 12, 31)

    query = f"""
    SELECT e.emp_no, e.gender, e.birth_date, e.hire_date, s.salary, s.from_date, s.to_date
    FROM employees e
    JOIN salaries s ON e.emp_no = s.emp_no
    WHERE s.from_date <= '{end_date}' AND s.to_date >= '{start_date}';
    """

    # Execute the query
    cursor.execute(query)
    # Fetch all the results
    data = cursor.fetchall()
    # Create a DataFrame from the results
    curr_df = pd.DataFrame(data,
                           columns=['emp_no', 'gender', 'birth_date', 'hire_date', 'salary', 'from_date', 'to_date'])

    # Ensure dates are in datetime format
    curr_df['from_date'] = pd.to_datetime(curr_df['from_date'], errors='coerce')
    curr_df['to_date'] = pd.to_datetime(curr_df['to_date'], errors='coerce')

    # Calculate weighted annual salary for each row
    annual_salaries = []
    for index, row in curr_df.iterrows():
        # Determine the period of the salary that falls within the year
        period_start = max(row['from_date'], start_date)
        period_end = min(row['to_date'], end_date)

        if period_start <= period_end:
            days_worked = (period_end - period_start).days + 1
        else:
            days_worked = 0

        total_days = days_in_year(year)
        weighted_salary = (days_worked / total_days) * row['salary']
        annual_salaries.append(weighted_salary)

    curr_df['annual_salary'] = annual_salaries

    # Aggregate annual salary by employee
    annual_salary_by_emp = curr_df.groupby('emp_no').agg({
        'gender': 'first',
        'annual_salary': 'sum'
    }).reset_index()

    # Add the year column
    annual_salary_by_emp['year'] = year

    # Append to the main DataFrame
    df_annsal = pd.concat([df_annsal, annual_salary_by_emp], ignore_index=True)

# Calculate the gender pay gap for each year
pay_gap_df = df_annsal.groupby(['year', 'gender'])['annual_salary'].mean().unstack().reset_index()

# Calculate the pay gap as a percentage of male earnings
pay_gap_df['pay_gap'] = ((pay_gap_df['M'] - pay_gap_df['F']) / pay_gap_df['M']) * 100

# Plot the pay gap evolution
plt.figure(figsize=(10, 6))
plt.plot(pay_gap_df['year'], pay_gap_df['pay_gap'], marker='o', linestyle='-', color='b')

# Adding titles and labels
plt.title('Gender Pay Gap Evolution Over Years')
plt.xlabel('Year')
plt.ylabel('Pay Gap (%)')
plt.grid(True)
plt.xticks(pay_gap_df['year'])  # Set x-ticks to be each year for better readability
plt.axhline(0, color='gray', linestyle='--', linewidth=0.5)  # Adding a horizontal line at y=0 for reference

# Add text labels for the pay gap figures
for i, row in pay_gap_df.iterrows():
    plt.text(row['year'], row['pay_gap'], f"{row['pay_gap']:.2f}%", fontsize=9, ha='right', va='bottom')

plt.savefig(f'{output_dir}/gender_paygap_over_time.png')
plt.close()

# 3.d. Generate slides
# Slide 3: The Gender Pay Gap
create_slide(prs, "ESRS Disclosure requirement S1-16: Gender Pay Gap",
             content=content)

# Slide 3: The Gender Pay Gap by Title
create_slide(prs, "ESRS Disclosure requirement S1-16: Gender Pay Gap by Title",
             img_path=f'{output_dir}\\gender_paygap_by_title.png')

# Slide 3: Net Evolution per Department
create_slide(prs, "ESRS Disclosure requirement S1-16: Gender Pay Gap by Department",
             img_path=f'{output_dir}\\gender_paygap_by_dept.png')

# Slide 3: The Gender Pay Gap over time
create_slide(prs, "ESRS Disclosure requirement S1-16: Gender Pay Gap (Annualized Average) over Time",
             img_path=f'{output_dir}\\gender_paygap_over_time.png')

# 4. Employees df
# 4.a. Fetching & processing data
# 4.a. 1.- Fetching employees df
# 4.a. 1.- (a) Employee Hire Data
query = """
SELECT e.emp_no, e.gender, e.hire_date, d.max_to_date
FROM employees e
JOIN (
    SELECT emp_no, MAX(to_date) AS max_to_date
    FROM dept_emp
    GROUP BY emp_no
) d ON e.emp_no = d.emp_no
ORDER BY e.hire_date;
"""
cursor.execute(query)
data = cursor.fetchall()
emp_df = pd.DataFrame(data, columns=['emp_no', 'gender', 'hire_date', "to_date"])
num_rows1 = emp_df.shape[0]
print(f'The number of rows in the dataset is: {num_rows1}')

# 4.a. 2.- (b) Fix last date
# Replace '9999-01-01' with a placeholder, such as pd.NaT (Not a Time)
emp_df['to_date'] = pd.to_datetime(emp_df['to_date'], errors='coerce')
emp_df['to_date'] = emp_df['to_date'].replace(pd.Timestamp('9999-01-01'), pd.NaT)

# Find the latest real date in the 'end_date' column
latest_real_date = emp_df['to_date'].max(skipna=True)

# Do not add two months to this latest date
new_date = latest_real_date

# Replace NaT (which were originally '9999-01-01') with the new calculated date
emp_df['to_date'] = emp_df['to_date'].fillna(new_date)

# Debug: Print first 15 rows of the merged DataFrame
print("Merged Employee Data (hire_date and to_date):")
print(emp_df.head(15))

# Check the data types of each column
print(emp_df.dtypes)

# hire_date is not yet defined as datetime.date
emp_df['hire_date'] = pd.to_datetime(emp_df['hire_date'])

# 4.a. 2.- (c) Generate df to plot evolution of employees
# Generate date range
date_range = pd.date_range(start=emp_df['hire_date'].min(), end=emp_df['to_date'].max(), freq='D')

# Initialize a DataFrame to hold the counts
employee_counts = pd.DataFrame(date_range, columns=['date'])
employee_counts['count'] = 0

# Update the count for each date
for _, row in emp_df.iterrows():
    employee_counts.loc[(employee_counts['date'] >= row['hire_date'])
                        &
                        (employee_counts['date'] <= row['to_date']), 'count'] += 1

# 4.a. 2.- (c) Generate df to plot evolution of employees
# Initialize DataFrames to hold the counts (we keep using the date_range
employee_counts_male = pd.DataFrame(date_range, columns=['date'])
employee_counts_male['count'] = 0

employee_counts_female = pd.DataFrame(date_range, columns=['date'])
employee_counts_female['count'] = 0

# Update the count for each date for male employees
for _, row in emp_df[emp_df['gender'] == 'M'].iterrows():
    employee_counts_male.loc[(employee_counts_male['date'] >= row['hire_date'])
                             &
                             (employee_counts_male['date'] <= row['to_date']), 'count'] += 1

# Update the count for each date for female employees
for _, row in emp_df[emp_df['gender'] == 'F'].iterrows():
    employee_counts_female.loc[(employee_counts_female['date'] >= row['hire_date'])
                               &
                               (employee_counts_female['date'] <= row['to_date']), 'count'] += 1

# 4.c. Generate plots
# 4.c. 1.- Plot 1: Change in number of employees
plt.figure(figsize=(10, 6))
plt.plot(employee_counts['date'], employee_counts['count'], label='Number of Employees')
plt.title('Number of Employees Over Time')
plt.xlabel('Date')
plt.ylabel('Number of Employees')
plt.grid(True)
plt.savefig(f'{output_dir}/num_employees_over_time.png')
plt.close()

# 3.c. 2.- Plot 2: Change in number of male and female employees
# Plotting
plt.figure(figsize=(10, 6))

# Plot for male employees
plt.plot(employee_counts_male['date'], employee_counts_male['count'], label='Male Employees', color='blue')

# Plot for female employees
plt.plot(employee_counts_female['date'], employee_counts_female['count'], label='Female Employees', color='red')

plt.xlabel('Date')
plt.ylabel('Number of Employees')
plt.title('Temporal Evolution of Total Number of Employees by Gender')
plt.legend()
plt.grid(True)
plt.savefig(f'{output_dir}/num_employees_per_gender_over_time.png')
plt.close()

# 4.d. Generate slides
# Slide 2: Line Plot
create_slide(prs, "Number of Employees Over Time",
             img_path=f'{output_dir}\\num_employees_over_time.png')

# Slide 3: Average Hired Each Month
create_slide(prs, "ESRS S1-7: Gender Workforce Composition Over Time",
             img_path=f'{output_dir}\\num_employees_per_gender_over_time.png')

# 5- Department employees df
# 5.a. Fetching data
# 5.a. 1.- Fetching department employees df
# 5.a. 1.- (a) Department Employee Data
query = """
SELECT from_date, to_date, dept_no, COUNT(emp_no) as num_employees
FROM dept_emp
GROUP BY from_date, to_date, dept_no
ORDER BY from_date;
"""
cursor.execute(query)
data = cursor.fetchall()
dept_emp_df = pd.DataFrame(data, columns=['from_date', 'to_date', 'dept_no', 'num_employees'])

# 5.a. 1.- (b) Fetching Department Data
query = """
SELECT dept_no, dept_name
FROM departments;
"""
cursor.execute(query)
data = cursor.fetchall()
dept_df = pd.DataFrame(data, columns=['dept_no', 'dept_name'])

# 5.b. 2.- Processing department employees df
# 5.b. 2.- (a) Merge department names into the department employee data
dept_emp_df = dept_emp_df.merge(dept_df, on='dept_no', how='left')

# Debug: Print first 15 rows of the merged DataFrame
print("Merged Employee Data (from_date and to_date):")
print(dept_emp_df.head(15))

# 5.b. 2.- (b) Fix last date
# Replace '9999-01-01' with a placeholder, such as pd.NaT (Not a Time)
dept_emp_df['to_date'] = pd.to_datetime(dept_emp_df['to_date'], errors='coerce')
dept_emp_df['to_date'] = dept_emp_df['to_date'].replace(pd.Timestamp('9999-01-01'), pd.NaT)

# Replace NaT (which were originally '9999-01-01') with the new (previously calculated) date
dept_emp_df['to_date'] = dept_emp_df['to_date'].fillna(new_date)

# Debug: Print first 15 rows of the merged DataFrame
print("Merged Employee Data (hire_date and to_date):")
print(dept_emp_df.head(15))

# Check the data types of each column
print(dept_emp_df.dtypes)

# hire_date is not yet defined as datetime.date
dept_emp_df['from_date'] = pd.to_datetime(dept_emp_df['from_date'])

# Generate date range
date_range = pd.date_range(start=dept_emp_df['from_date'].min(), end=dept_emp_df['to_date'].max(), freq='D')

# Get unique department names
departments = dept_emp_df['dept_name'].unique()

# Initialize a dictionary to hold DataFrames for each department
dept_counts = {}
for dept in departments:
    dept_counts[dept] = pd.DataFrame(date_range, columns=['date'])
    dept_counts[dept]['count'] = 0

# Update the count for each date for each department
for dept in departments:
    for _, row in dept_emp_df[dept_emp_df['dept_name'] == dept].iterrows():
        dept_counts[dept].loc[(dept_counts[dept]['date'] >= row['from_date'])
                              &
                              (dept_counts[dept]['date'] <= row['to_date']), 'count'] += 1

# 4.c. Generate Plots
# 4.c. 1.- Plot 3: Change in number of employees per department
plt.figure(figsize=(10, 6))

# Plot for each department
for dept in departments:
    plt.plot(dept_counts[dept]['date'], dept_counts[dept]['count'], label=f'{dept} Employees')

plt.xlabel('Date')
plt.ylabel('Number of Employees')
plt.title('Evolution of Total Number of Employees by Department')
plt.legend()
plt.grid(True)
plt.savefig(f'{output_dir}/num_employees_per_department_over_time.png')
plt.close()

# 4.d. Generate slides
# Slide 3: Net Evolution per Department
create_slide(prs, "Workforce Composition by Department Over Time",
             img_path=f'{output_dir}\\num_employees_per_department_over_time.png')

# 5. Salaries df
# 5.a. Fetching data
# 5.a. 1.- Fetching Employee Hire Data
# We use the hired_df

# 5.a. 2.- Fetching Employee Salary Data
query = """
SELECT emp_no, salary, from_date, to_date
FROM salaries;
"""
cursor.execute(query)
data = cursor.fetchall()
sal_df = pd.DataFrame(data, columns=['emp_no', 'salary', 'from_date', 'to_date'])

# 5.b. Processing employee salaries df
# 5.b. 2.- Merge employees df into the employee salaries data
# Change to_date in emp_df
emp_df.rename(columns={'to_date': 'leave_date'}, inplace=True)

# Merge
sal_emp_df = sal_df.merge(emp_df, on='emp_no', how='left')

# Debug: Print first 15 rows of the merged DataFrame
print("Merged Employee Salaries Data (from_date and to_date):")
print(sal_emp_df.head(15))

# 5.b. 3.- Fix last date
# Replace '9999-01-01' with a placeholder, such as pd.NaT (Not a Time)
sal_emp_df['to_date'] = pd.to_datetime(sal_emp_df['to_date'], errors='coerce')
sal_emp_df['to_date'] = dept_emp_df['to_date'].replace(pd.Timestamp('9999-01-01'), pd.NaT)

# Replace NaT (which were originally '9999-01-01') with the new (previously calculated) date
sal_emp_df['to_date'] = sal_emp_df['to_date'].fillna(new_date)

# Debug: Print first 15 rows of the merged DataFrame
print("Merged Employee Salaries Data (hire_date and to_date):")
print(sal_emp_df.head(15))

# 5.c. Analysing department employees df
# 5.c. 1.- Calculate average starting salary for men and women
# Identify starting salary
df_starting = sal_emp_df.loc[sal_emp_df.groupby('emp_no')['from_date'].idxmin()]

# Calculate average starting salary
avg_starting_salary = df_starting.groupby('gender')['salary'].mean()
print(avg_starting_salary)

# 5.c. 2.- Analyze Salary Progression
# Create a regression model to analyze salary progression controlling for starting salary and tenure
# Check the data types of each column
print(sal_emp_df.dtypes)

# Adding tenure
sal_emp_df['tenure_days'] = (sal_emp_df['leave_date'] - sal_emp_df['hire_date']).dt.days

# Map gender to numeric values
sal_emp_df['gender_dum'] = sal_emp_df['gender'].map({'M': 0, 'F': 1})

# Create the interaction term
sal_emp_df['interaction'] = sal_emp_df['tenure_days'] * sal_emp_df['gender_dum']

# Initialize regression model
X = sal_emp_df[['tenure_days', 'gender_dum', 'interaction']]
X = sm.add_constant(X)
y = sal_emp_df['salary']
model = sm.OLS(y, X).fit()

# Print regression summary
print(model.summary())

# 5.c. 3.- Analyze evolution in starting salaries
# Create a new column for year of hire
sal_emp_df['hire_year'] = sal_emp_df['hire_date'].dt.year

# Group by hire year and gender, then calculate the average starting salary
df_starting_year = df_starting.copy()
df_starting_year['hire_year'] = df_starting_year['hire_date'].dt.year
avg_starting_salary_year = df_starting_year.groupby(['hire_year', 'gender'])['salary'].mean().unstack()
'''
# 5.c. 3.- Analyze evolution in starting salaries
# 5.c. 3.- (a)
# Basic multilevel model: Salary ~ Tenure + (1 | emp_no)
model2 = mixedlm("salary ~ tenure_days", sal_emp_df, groups=sal_emp_df["emp_no"])
result2 = model2.fit()
print(result2.summary())

# Extended model with interaction: Salary ~ Tenure * Gender + (1 | emp_no)
model_interaction2 = mixedlm("salary ~ tenure_days * gender_dum", sal_emp_df, groups=sal_emp_df["emp_no"])
result_interaction2 = model_interaction2.fit()
print(result_interaction2.summary())
'''
# 5.d. Plot data
# 5.d. 1.- Plot salary progression
# Create scatter plot
sns.scatterplot(data=sal_emp_df, x='tenure_days', y='salary', hue='gender')

# Calculate regression lines
tenure_range = np.linspace(sal_emp_df['tenure_days'].min(), sal_emp_df['tenure_days'].max(), 100)
male_slope = model.params['tenure_days']
female_slope = model.params['tenure_days'] + model.params['interaction']
male_intercept = model.params['const']
female_intercept = model.params['const'] + model.params['gender_dum']

plt.figure(figsize=(10, 6))

# Male regression line
plt.plot(tenure_range, male_slope * tenure_range + male_intercept, color='blue', label='Male Regression Line')

# Female regression line
plt.plot(tenure_range, female_slope * tenure_range + female_intercept, color='red', label='Female Regression Line')

# Calculate average salary for each tenure_days value
avg_salary_by_tenure = sal_emp_df.groupby('tenure_days')['salary'].mean().reset_index()

# Plot average salary
plt.scatter(avg_salary_by_tenure['tenure_days'], avg_salary_by_tenure['salary'], color='#F8E7FF', label='Average Salary',
            marker='.')

# Customize the plot
plt.xlabel('Tenure Days')
plt.ylabel('Salary')
plt.title('Interaction Effect of Tenure and Gender on Salary')
plt.legend()
plt.grid(True)
plt.savefig(f'{output_dir}/sal_prog_by_gender_over_time.png')
plt.close()

# 5.d. 2.- Plot the historical perspective of starting salaries
# Define the figure with a wider aspect ratio and adjust margins
fig, ax = plt.subplots(figsize=(16, 9))  # Wider and slightly taller
avg_starting_salary_year.plot(kind='line', ax=ax)# Plot the data
ax.set_xlabel('Year of Hire') # Adding labels and title
ax.set_ylabel('Average Starting Salary')
ax.set_title('Historical Perspective of Starting Salaries by Gender')
ax.legend() # Adding legend and grid
ax.grid(True)

plt.tight_layout() # Adjusting layout to prevent clipping
plt.savefig(f'{output_dir}/start_sal_by_gender_over_time.png', dpi=300) # Saving the plot to a file with higher
# DPI for better resolution
plt.close() # Closing the plot
# For salary progression, we need a more complex regression analysis over time

# 5.d. Generate slides
# Slide 4: Salary Progression by Gender
create_slide(prs, "Salary Progression by Gender (Average Across the Company Lifetime)",
             img_path=f'{output_dir}\\sal_prog_by_gender_over_time.png')

# Slide 5: Starting Salary Evolution by Gender
create_slide(prs, "Average Starting Salary by Gender Over Time",
             img_path=f'{output_dir}\\start_sal_by_gender_over_time.png')

# Save the presentation
prs.save(f'{output_dir}\\Employees_Report.pptx')
